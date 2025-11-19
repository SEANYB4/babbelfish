import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from pptx import Presentation
from googletrans import Translator
from pptx.util import Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageTk
import time



# RESIZE LOGO IMAGE



# Extract text from PPTX file
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run is not None:
                        text_runs.append(run.text)

    return text_runs

# Translate text
def translate_text(texts, dest='zh-cn', progress_callback=None):
    translator = Translator()
    translations = []
    for index, text in enumerate(texts):
        try:
            translation= translator.translate(text, dest=dest)
            translations.append(translation.text)

        except Exception as e:
            translations.append('Translation error')
            print(f"Error: {e}")

        if progress_callback:
            progress_callback(index + 1, len(texts))
    return translations


# Append or replace text in PPTX
def modify_pptx(file_path, original_texts, translations, method='append', progress_callback=None, save_path="translated_presentation.pptx"):
    prs = Presentation(file_path)
    i = 0
    # num_texts = len(translations)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if i < len(translations) and run.text in original_texts:
                        if method == 'replace':
                            run.text = translations[i]
                        elif method == 'append':
                            run.text += " " + translations[i]
                        i += 1

                        # if progress_callback:
                        #     progress_callback(i, num_texts)

    prs.save(save_path)

class App(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("PPT-Translator")
        

        # Open image with Pillow
        self.original_image = Image.open('LOGO.png')
        self.logo_image = ImageTk.PhotoImage(self.original_image)

        # Set the window size to the image size
        # width, height = self.original_image.size
        # self.geometry(f"{width}x{height}")

        
        

        # Logo
        
        self.logo_label = ttk.Label(self, image=self.logo_image)
        self.logo_label.place(relwidth=1, relheight=1)

        self.geometry("600x400")

         # Navigation Frame
        nav_frame = tk.Frame(self, height=100, bg='#ADADAD')
        nav_frame.grid(column=0, row=0, sticky='ew')


        self.grid_columnconfigure(0, weight=1)

        nav_frame.columnconfigure(0, weight=1)
        nav_frame.columnconfigure(1, weight=1)
        nav_frame.columnconfigure(2, weight=1)
        nav_frame.columnconfigure(3, weight=1)


        # Styling for ttk Button
        style = ttk.Style()
        style.configure("LightBlue.TButton", foreground='black', background='#ADD8E6')
        style.configure("Status.TLabel", foreground='black', background='#D3FDAD', font=('Hevetica', 10))
       

        # Widgets
        self.file_label = ttk.Label(nav_frame, text="Select ppt file: ")
        self.file_label.grid(column=0, row=0, pady=5, sticky='ew')


        # Progress Bar
        self.progress = ttk.Progressbar(nav_frame, orient='horizontal', length=500, mode='determinate')
        self.progress.grid(row=3, column=0, columnspan=3, sticky='ew')


        # Open PPT Button
        self.open_file_button = ttk.Button(nav_frame, text="Open PPTX File", command = self.load_file, style="LightBlue.TButton").grid(row=0, column=2, sticky='ew')


        self.language_label = ttk.Label(nav_frame, text="Select target language: ")
        self.language_label.grid(column=0, row=1, pady=5, sticky='ew')

        self.language_var = tk.StringVar()
        self.language_entry = ttk.Entry(nav_frame, textvariable=self.language_var)

        self.language_var.set("zh-cn")
        self.language_entry.grid(row=1, column=1, sticky='ew')


        self.translate_button = ttk.Button(nav_frame, text="Translate PPT", command=self.translate, style="LightBlue.TButton")
        self.translate_button.grid(row=1, column=2, pady=5, sticky='ew')


       

        self.status_label = ttk.Label(nav_frame, text="Select a .pptx file to translate", style="Status.TLabel")
        self.status_label.grid(row=0, column=1, pady=5, sticky='ew')

       
        self.info_label = ttk.Label(self, text="REQUIRES INTERNET CONNECTION TO CONNECT WITH TRANSLATION SERVICE", background='#D3FDAD', font=('Helvetica', 10))
        self.info_label.grid(column=0, row=1, sticky='ew')

        self.info_label2 = ttk.Label(self, text="REQUIRES VPN ENABLED IF INSIDE CHINA", background='#D3FDAD', font=('Helvetica', 10))
        self.info_label2.grid(column=0, row=2, sticky='ew')


        # Data
        self.file_path = None
        self.original_texts = []
        self.translations = []


        self.last_resize = time.time()



        self.after(1, self.auto_resize_window)
        self.after(1000, self.resize_image(None))

        self.resizable(False, False)
        # self.bind("<Configure>", self.resize_image)
        
       
            

    def update_progress(self, current, total):
        self.progress['value'] = (current / total) * 100
        self.update_idletasks()


    def load_file(self):
        self.file_path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])

        if self.file_path:
            # self.file_label.config(text=f"Selected File: {self.file_path}")
            self.original_texts = extract_text_from_pptx(self.file_path)
            self.status_label.config(text=f"Selected File: {self.file_path}")
            self.progress['value'] = 0
            self.progress['maximum'] = 100
            self.update_idletasks()

    def translate(self):
        if not self.original_texts:
            messagebox.showerror("Error", "No text extracted or file not loaded.")
            return
        
        self.progress['value'] = 0
        self.progress['maximum'] = 100
        self.translations = translate_text(self.original_texts, self.language_var.get(), self.update_progress)
        self.status_label.config(text="Translation complete. Saving translated .pptx file.")
        self.save()
    
    def save(self):
        if not self.translations:
            messagebox.showerror("Error", "No translations available.")
            return
        
        save_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint Presentation", "*.pptx")], title="Save Translated Presentation As")

        if not save_path:
            messagebox.showinfo("Cancelled", "Save operation cancelled.")


        self.progress['value'] = 0
        self.progress['maximum'] = 100
        self.update_idletasks()
        modify_pptx(self.file_path, self.original_texts, self.translations, method='append', progress_callback=self.update_progress, save_path=save_path)

        messagebox.showinfo("Success", f"Presentation saved as '{save_path}'")

        
    def resize_image(self, event):

        if time.time() - self.last_resize<0.5:
            return
        
        self.last_resize = time.time()


    # Check if the window is minimized to an extent where width or height becomes zero
        new_width = self.logo_label.winfo_width()
        new_height = self.logo_label.winfo_height()

        if new_width > 0 and new_height > 0:  # Ensure dimensions are greater than zero
            # Resize the original image using Image.ANTIALIAS for better quality
            resized_image = self.original_image.resize((new_width, new_height), Image.LANCZOS)
            self.logo_image = ImageTk.PhotoImage(resized_image)

            # Update the label's image
            self.logo_label.configure(image=self.logo_image)
            self.logo_label.image = self.logo_image  # Keep a reference to avoid garbage collection
        else:
            # Handle the case where dimensions are zero (e.g., window is minimized)
            print("Window is too small for image resizing.")


    def auto_resize_window(self):

        new_width = 1000
        new_height = 600

        self.geometry(f"{new_width}x{new_height}")

        resized_image = self.original_image.resize((new_width, new_height), Image.LANCZOS)
        self.logo_image = ImageTk.PhotoImage(resized_image)
        self.logo_label.configure(image=self.logo_image)
        self.logo_label.image = self.logo_image  # Maintain reference to avoid garbage collection






if __name__ == "__main__":
    app = App()
    app.mainloop()